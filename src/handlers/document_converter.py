import io
import os
import logging
import tempfile
from pathlib import Path
from telegram import Update
from telegram.ext import ContextTypes
from docx2pdf import convert
from pdf2docx import Converter
from PyPDF2 import PdfReader, PdfWriter
import mammoth
from src.keyboards import get_doc_format_keyboard
from src.config import IMAGES
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

async def convert_docx_to_pptx(docx_path: str, pptx_path: str) -> bool:
    """Convert DOCX document to PPTX presentation."""
    try:
        # Открываем документ Word
        doc = Document(docx_path)
        # Создаем новую презентацию
        prs = Presentation()
        
        # Настройки слайда
        title_slide_layout = prs.slide_layouts[0]  # Title Slide
        content_slide_layout = prs.slide_layouts[1]  # Title and Content
        
        # Создаем титульный слайд
        title_slide = prs.slides.add_slide(title_slide_layout)
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        # Устанавливаем заголовок из первого параграфа документа
        if doc.paragraphs:
            title.text = doc.paragraphs[0].text
            if len(doc.paragraphs) > 1:
                subtitle.text = doc.paragraphs[1].text
        
        # Обрабатываем остальной текст
        current_slide = None
        current_text = []
        
        for para in doc.paragraphs[2:]:
            # Если это заголовок, создаем новый слайд
            if para.style.name.startswith('Heading'):
                # Сохраняем текст предыдущего слайда
                if current_slide and current_text:
                    content = current_slide.shapes.placeholders[1]
                    content.text = "\n".join(current_text)
                    current_text = []
                
                # Создаем новый слайд
                current_slide = prs.slides.add_slide(content_slide_layout)
                title = current_slide.shapes.title
                title.text = para.text
            else:
                # Добавляем текст к текущему слайду
                if para.text.strip():
                    current_text.append(para.text)
                    
                    # Если накопилось много текста, создаем новый слайд
                    if len(current_text) >= 5:
                        if not current_slide:
                            current_slide = prs.slides.add_slide(content_slide_layout)
                            title = current_slide.shapes.title
                            title.text = "Продолжение"
                        
                        content = current_slide.shapes.placeholders[1]
                        content.text = "\n".join(current_text)
                        current_text = []
                        current_slide = None
        
        # Сохраняем последний слайд
        if current_slide and current_text:
            content = current_slide.shapes.placeholders[1]
            content.text = "\n".join(current_text)
        
        # Сохраняем презентацию
        prs.save(pptx_path)
        return True
        
    except Exception as e:
        logger.error(f"Error converting DOCX to PPTX: {str(e)}", exc_info=True)
        return False

async def handle_document_conversion(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle document messages for conversion."""
    try:
        document = update.message.document
        mime_type = document.mime_type
        file_name = document.file_name
        
        supported_types = {
            'application/pdf': 'PDF',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'DOCX',
            'application/msword': 'DOC',
            'text/plain': 'TXT'
        }
        
        if mime_type not in supported_types:
            await update.message.reply_text(
                text=f"{IMAGES['error']} Извините, этот формат не поддерживается.\n\n"
                     "Поддерживаемые форматы:\n"
                     "• PDF\n"
                     "• DOCX\n"
                     "• DOC\n"
                     "• TXT"
            )
            return
            
        # Download document
        doc_file = await context.bot.get_file(document.file_id)
        doc_bytes = await doc_file.download_as_bytearray()
        
        # Store document info in context
        context.user_data['current_document'] = {
            'bytes': doc_bytes,
            'name': file_name,
            'type': supported_types[mime_type]
        }
        
        # Если это DOCX, предлагаем дополнительные опции
        if supported_types[mime_type] == 'DOCX':
            await update.message.reply_text(
                text=f"{IMAGES['formats']}\nВыберите формат для конвертации:\n"
                     "1. PDF\n"
                     "2. PPTX (Презентация)\n"
                     "Отправьте номер нужного формата (1 или 2):"
            )
            # Сохраняем путь к файлу для последующей конвертации
            with tempfile.TemporaryDirectory() as temp_dir:
                source_path = os.path.join(temp_dir, file_name)
                with open(source_path, 'wb') as f:
                    f.write(doc_bytes)
                context.user_data['current_file'] = {
                    'path': source_path,
                    'name': file_name,
                    'mime_type': mime_type
                }
        else:
            # Для других форматов показываем стандартную клавиатуру
            keyboard = get_doc_format_keyboard(supported_types[mime_type])
            await update.message.reply_text(
                text=f"{IMAGES['formats']}\nВыберите формат для конвертации:",
                reply_markup=keyboard
            )
        
    except Exception as e:
        logger.error(f"Error handling document: {str(e)}")
        await update.message.reply_text(
            text=f"{IMAGES['error']} Извините, произошла ошибка при обработке документа. Пожалуйста, попробуйте снова."
        )

async def convert_document(update: Update, context: ContextTypes.DEFAULT_TYPE, target_format: str) -> None:
    """Convert document to target format."""
    try:
        if 'current_document' not in context.user_data:
            await update.message.reply_text(
                text=f"{IMAGES['error']} Пожалуйста, сначала отправьте документ для конвертации."
            )
            return
            
        doc_info = context.user_data['current_document']
        source_bytes = doc_info['bytes']
        source_type = doc_info['type']
        original_name = Path(doc_info['name']).stem
        
        # Create temporary directory for conversion
        with tempfile.TemporaryDirectory() as temp_dir:
            source_path = os.path.join(temp_dir, f"source.{source_type.lower()}")
            target_path = os.path.join(temp_dir, f"target.{target_format.lower()}")
            
            # Write source file
            with open(source_path, 'wb') as f:
                f.write(source_bytes)
            
            # Convert based on source and target formats
            if source_type == 'PDF' and target_format == 'DOCX':
                cv = Converter(source_path)
                cv.convert(target_path)
                cv.close()
            
            elif source_type in ['DOC', 'DOCX'] and target_format == 'PDF':
                convert(source_path, target_path)
            
            elif source_type in ['DOC', 'DOCX'] and target_format == 'PPTX':
                success = await convert_docx_to_pptx(source_path, target_path)
                if not success:
                    raise Exception("Failed to convert to PPTX")
            
            elif source_type in ['DOC', 'DOCX'] and target_format == 'TXT':
                with open(source_path, 'rb') as docx_file:
                    result = mammoth.extract_raw_text(docx_file)
                    with open(target_path, 'w', encoding='utf-8') as txt_file:
                        txt_file.write(result.value)
            
            elif source_type == 'TXT' and target_format in ['PDF', 'DOCX']:
                doc = Document()
                with open(source_path, 'r', encoding='utf-8') as txt_file:
                    content = txt_file.read()
                    doc.add_paragraph(content)
                doc.save(target_path if target_format == 'DOCX' else source_path + '.docx')
                
                if target_format == 'PDF':
                    convert(source_path + '.docx', target_path)
            
            # Read converted file
            with open(target_path, 'rb') as f:
                output_bytes = f.read()
            
        # Send converted file with appropriate caption
        caption = "Вот ваш документ в формате"
        if target_format == 'PPTX':
            caption = f"{caption} PPTX! ✨\n• Заголовки документа преобразованы в слайды\n• Текст разделен на удобные для чтения части\n• Сохранено форматирование заголовков"
        else:
            caption = f"{caption} {target_format}! ✨"
            
        await update.message.reply_document(
            document=output_bytes,
            filename=f"{original_name}.{target_format.lower()}",
            caption=caption
        )
        
        # Send success message
        success_msg = f"{IMAGES['success']} Конвертация завершена успешно!\n\n📄 Исходный формат: {source_type}\n📑 Новый формат: {target_format}"
        if target_format == 'PPTX':
            success_msg += "\n\n💡 Презентация создана на основе структуры документа"
            
        await update.message.reply_text(text=success_msg)
        
        # Clear the stored document
        del context.user_data['current_document']
        
    except Exception as e:
        logger.error(f"Error converting document: {str(e)}")
        await update.message.reply_text(
            text=f"{IMAGES['error']} Извините, произошла ошибка при конвертации. Пожалуйста, попробуйте снова."
        )

async def handle_conversion_choice(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle user's choice for document conversion."""
    try:
        if 'current_file' not in context.user_data:
            await update.message.reply_text(
                text=f"{IMAGES['error']} Пожалуйста, сначала отправьте документ."
            )
            return
        
        choice = update.message.text.strip()
        file_info = context.user_data['current_file']
        source_path = file_info['path']
        file_name = file_info['name']
        
        with tempfile.TemporaryDirectory() as temp_dir:
            if choice == '1':
                # Convert DOCX to PDF
                target_path = os.path.join(temp_dir, Path(file_name).stem + '.pdf')
                convert(source_path, target_path)
                
                # Read converted file
                with open(target_path, 'rb') as f:
                    output_bytes = f.read()
                
                # Send converted file
                await update.message.reply_document(
                    document=output_bytes,
                    filename=Path(file_name).stem + '.pdf',
                    caption="Вот ваш документ в формате PDF! ✨"
                )
                
            elif choice == '2':
                # Convert DOCX to PPTX
                target_path = os.path.join(temp_dir, Path(file_name).stem + '.pptx')
                success = await convert_docx_to_pptx(source_path, target_path)
                
                if success:
                    # Read converted file
                    with open(target_path, 'rb') as f:
                        output_bytes = f.read()
                    
                    # Send converted file
                    await update.message.reply_document(
                        document=output_bytes,
                        filename=Path(file_name).stem + '.pptx',
                        caption="Вот ваша презентация! ✨\n"
                                "• Заголовки документа преобразованы в слайды\n"
                                "• Текст разделен на удобные для чтения части\n"
                                "• Сохранено форматирование заголовков"
                    )
                else:
                    await update.message.reply_text(
                        text=f"{IMAGES['error']} Не удалось конвертировать документ в презентацию."
                    )
            
            else:
                await update.message.reply_text(
                    text=f"{IMAGES['error']} Пожалуйста, выберите 1 (PDF) или 2 (PPTX)."
                )
        
        # Clear stored file info
        del context.user_data['current_file']
        
    except Exception as e:
        logger.error(f"Error handling conversion choice: {str(e)}", exc_info=True)
        await update.message.reply_text(
            text=f"{IMAGES['error']} Произошла ошибка при конвертации документа."
        ) 